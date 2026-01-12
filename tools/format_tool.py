import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches

def create_pdf(text: str, filename: str):
    """
    Creates a simple PDF from text.
    """
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    
    # Simple line handling
    for line in text.split('\n'):
        # Encode to latin-1 to avoid unicode errors in standard fpdf, or replace characters
        safe_line = line.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 10, safe_line)
        
    pdf.output(filename)
    return filename

def create_ppt(text: str, filename: str):
    """
    Creates a simple PPT from text. 
    Assumes text is somewhat structured (bullet points).
    """
    prs = Presentation()
    slide_layout = prs.slide_layouts[1] # Bullet slide
    
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Research Report"
    
    tf = body_shape.text_frame
    tf.text = text[:1000] # Truncate for slide to avoid overflow
    
    prs.save(filename)
    return filename

from docx import Document

def create_markdown(text: str, filename: str):
    with open(filename, 'w') as f:
        f.write(text)
    return filename

def create_docx(text: str, filename: str):
    """
    Creates a simple Word Doc from text.
    """
    doc = Document()
    doc.add_heading('Research Report', 0)
    
    for line in text.split('\n'):
        doc.add_paragraph(line)
        
    doc.save(filename)
    return filename
